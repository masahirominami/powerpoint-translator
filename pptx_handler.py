# pptx_handler.py
from pptx import Presentation
from copy import deepcopy

from pptx import Presentation

# extract_text_from_pptx extracts all text from a PowerPoint file
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    collected = []

    for idx, slide in enumerate(prs.slides, start=1):
        collected.append(f"--- Slide {idx} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                collected.append(shape.text.strip())
        collected.append("")  # Add spacing between slides

    return "\n".join(collected)

# translate_text_in_place translates text in a PowerPoint file and saves it to
# a new file
def translate_text_in_place(pptx_path, output_path, target_lang):
    from translator import translate_text
    prs = Presentation(pptx_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text
                        translated_text = translate_text(original_text, target_lang)
                        run.text = translated_text  # ✅ preserves font style


    prs.save(output_path)
    print(f"✔️ Translated (in-place) file saved to: {output_path}")

